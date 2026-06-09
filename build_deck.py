"""Build the ALCF presentation from the template (v2 — addresses QA issues).

Key fixes vs v1:
- 4-block "big idea" layouts: split header (large bold) from body (smaller, normal)
  using two paragraphs with explicit font size on the body paragraph.
- All subtitle text constrained to a single line (~80 chars).
- All bullet lists trimmed so the body placeholder does not overflow.
- Title slide image placeholder removed (no image to insert).
- Quoted text cleaned (typos fixed only when clearly a typo).
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
TEMPLATE = os.path.join(HERE, "ALCF Presentation Template.pptx")
OUT = os.path.join(HERE, "ai-agent-coding-alcf.pptx")

# ============================================================================
# Helpers
# ============================================================================

def remove_all_slides(pres):
    sldIdLst = pres.slides._sldIdLst
    part = pres.part
    rid_slide = [(sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"), sldId)
                 for sldId in list(sldIdLst)]
    for rId, sldId in rid_slide:
        sldIdLst.remove(sldId)
        part.drop_rel(rId)

def remove_placeholder(slide, idx):
    """Remove a placeholder from a slide entirely (so the layout's empty
    'click to insert image' prompt doesn't render)."""
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == idx:
            sp = ph._element
            sp.getparent().remove(sp)
            return True
    return False

def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    raise KeyError(f"no placeholder idx={idx} in slide; have {[p.placeholder_format.idx for p in slide.placeholders]}")

def set_text(ph, text, *, size=None, bold=None):
    """Single-paragraph text into a placeholder, optionally with size/bold override."""
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    run = p0.add_run()
    run.text = text
    if size is not None: run.font.size = Pt(size)
    if bold is not None: run.font.bold = bold

def set_block(ph, header, body, *, body_size=14):
    """Big-idea block: header (default style = large bold) + smaller body paragraph.

    The 4-block layout's level-1 default is 32pt bold white. We keep that for the
    header. For the body, we add a second paragraph and explicitly drop to 14pt
    normal weight (color is inherited as scheme bg1 = white)."""
    tf = ph.text_frame
    # clear
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    # Header — inherits 32pt bold from level 1 default
    run_h = p0.add_run()
    run_h.text = header
    # Body paragraph at smaller size
    p1 = tf.add_paragraph()
    run_b = p1.add_run()
    run_b.text = body
    run_b.font.size = Pt(body_size)
    run_b.font.bold = False
    # Force color = white (since we overrode bold/size, color inheritance can be
    # unreliable across renderers; explicit is safer).
    run_b.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

def set_bullets(ph, items):
    """Populate a content placeholder with bullet items (each either text or (text, level))."""
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    paragraphs = [p0]
    for _ in range(len(items) - 1):
        paragraphs.append(tf.add_paragraph())
    for item, para in zip(items, paragraphs):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        para.level = level
        run = para.add_run()
        run.text = text

def set_code(ph, code, *, font="Consolas", size=9):
    """Populate a placeholder with monospaced code, no bullets, no auto-indent.

    Overrides the master's spcBef (space before paragraph) and lnSpc (line spacing)
    so consecutive code lines pack tightly rather than spreading like bullets.
    """
    from pptx.oxml.ns import qn
    from lxml.etree import SubElement
    tf = ph.text_frame
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    for para in list(tf.paragraphs[1:]):
        para._p.getparent().remove(para._p)
    lines = code.split("\n")
    paragraphs = [p0]
    for _ in range(len(lines) - 1):
        paragraphs.append(tf.add_paragraph())
    for line, para in zip(lines, paragraphs):
        pPr = para._p.get_or_add_pPr()
        # Remove existing bullets, line-spacing, and space-before overrides
        for tag in ("a:buChar", "a:buAutoNum", "a:buNone", "a:lnSpc", "a:spcBef", "a:spcAft"):
            for child in pPr.findall(qn(tag)):
                pPr.remove(child)
        # Tight line spacing (100%) and zero space-before/after
        lnSpc = SubElement(pPr, qn("a:lnSpc"))
        SubElement(lnSpc, qn("a:spcPct")).set("val", "100000")
        spcBef = SubElement(pPr, qn("a:spcBef"))
        SubElement(spcBef, qn("a:spcPts")).set("val", "0")
        spcAft = SubElement(pPr, qn("a:spcAft"))
        SubElement(spcAft, qn("a:spcPts")).set("val", "0")
        SubElement(pPr, qn("a:buNone"))
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        para.level = 0
        run = para.add_run()
        run.text = line if line else " "
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = False
        # Inherit color (master defines body text color); leave alone.

# ============================================================================
# Build
# ============================================================================

pres = Presentation(TEMPLATE)
remove_all_slides(pres)
LAYOUTS = pres.slide_layouts
def add(layout_idx):
    return pres.slides.add_slide(LAYOUTS[layout_idx])

# --- 1. Title -------------------------------------------------------------
s = add(0)
set_text(get_ph(s, 0), "From Directing to Dialogue: Fifteen Weeks of AI Agent Coding for Performance Engineering")
set_text(get_ph(s, 1), "17 sessions · 2,650 prompts · 15 weeks")
set_text(get_ph(s, 17), "Brice Videau")
set_text(get_ph(s, 18), "Argonne Leadership Computing Facility\nArgonne National Laboratory\nJune 2026")
# Remove the empty image placeholder and unused presenter slots
for idx in (10, 19, 20, 21, 22):
    remove_placeholder(s, idx)

# --- 2. Agenda ------------------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "Agenda")
set_text(get_ph(s, 13), "Thirty minutes, two case studies, one tool")
set_bullets(get_ph(s, 14), [
    "Setup: how Claude Code was used at ALCF",
    "Project 1 — CCS: directing the agent on familiar code",
    "Project 2 — rust-gpu + claspr: collaborating across a knowledge gap",
    "Cross-cutting analysis: tokens, tools, sub-agents, context, memory",
    "Evolution of interaction style and tooling adoption",
    "Lessons learned and recommendations for ALCF performance engineering",
])

# --- 3. Why this evaluation (4-block) -------------------------------------
s = add(9)
set_text(get_ph(s, 0), "Why this evaluation")
set_text(get_ph(s, 13), "Coding agents are real; perf-eng needs to decide where they fit")
set_block(get_ph(s, 16), "Perf-eng work at ALCF",
    "Library and runtime work, kernel ports, profiling tooling, long-lived codebases, "
    "heavy review burden, frequent context switching.")
set_block(get_ph(s, 17), "Hypothesis",
    "An agent can carry implementation while the engineer directs design and review.")
set_block(get_ph(s, 18), "Success criteria",
    "Shippable PRs on real upstream repos. No regressions on familiar code. "
    "Credible progress on unfamiliar code. Defensible cost.")
set_block(get_ph(s, 19), "What this talk is not",
    "Not a benchmark and not a product comparison. A first-person retrospective from two ALCF-adjacent projects.")

# --- 4. Setup -------------------------------------------------------------
s = add(6)
set_text(get_ph(s, 0), "Setup: how Claude Code was used")
set_text(get_ph(s, 13), "Same tool, two laptops, one persona")
set_text(get_ph(s, 16), "Environment")
set_bullets(get_ph(s, 14), [
    "Two Ubuntu laptops: native Intel + Arm Mac (Ubuntu env)",
    "Claude Code CLI (Opus 4.x / Sonnet 4.x)",
    "Local toolchains: autotools, cargo, pocl, rusticl",
    "`gh` CLI for GitHub PRs, CI (incl. macOS runners)",
])
set_text(get_ph(s, 17), "Persona and isolation")
set_bullets(get_ph(s, 15), [
    "Dedicated GitHub identity: bricevideau-ai",
    "Dedicated gmail for notifications and 2FA",
    "Commits co-authored by Claude",
    "Per-project sessions, no cross-project mixing",
])
set_text(get_ph(s, 19), "What the agent keeps")
set_bullets(get_ph(s, 18), [
    "Per-project CLAUDE.md (orientation)",
    "Memory entries: feedback / project / reference",
    "Cross-session prompt history",
    "Sub-agent transcripts on disk",
])

# --- 5. Two case studies side by side -------------------------------------
s = add(5)
set_text(get_ph(s, 0), "The two case studies, side by side")
set_text(get_ph(s, 13), "Same tool, two very different collaboration modes")
set_text(get_ph(s, 16), "CCS (C Configuration Space)")
set_bullets(get_ph(s, 14), [
    "Autotuning configuration-space library",
    "C99, ~52K LOC, Python + Ruby bindings",
    "I wrote it and maintain it — high expertise",
    "Feb 25 → Apr 7 (~6 weeks, 4 sessions across both laptops)",
    "600 user prompts, 99 PRs (96 merged)",
    "Mode: direct, correct, ship",
])
set_text(get_ph(s, 17), "rust-gpu + claspr")
set_bullets(get_ph(s, 15), [
    "OpenCL Kernel target for rust-gpu; claspr single-source layer on top",
    "Rust + SPIR-V; ~142K + ~81K LOC",
    "Not a Rust dev; not a SPIR-V expert",
    "Mar 27 → Jun 8 (~10 weeks, 13 sessions across both laptops)",
    "2,059 user prompts, 24,540-line PR + greenfield repo",
    "Mode: explore, dialogue, decide, ship",
])

# --- 6. Session chronology ------------------------------------------------
# 17 sessions across two laptops; max parallelism is a 26-day overlap on
# rust-gpu (May 14 → Jun 8, afa1ce4c on Linux while 11e6e374 ran on the Mac).
s = add(3)  # Title, Subtitle and Bullets — gives a tall body for the timeline
set_text(get_ph(s, 0), "Session chronology")
set_text(get_ph(s, 13), "17 sessions across two laptops; up to 26 days running in parallel")
set_code(get_ph(s, 14),
"""                       Feb 25  Mar 09  Mar 27  Apr 07  Apr 24  May 04  May 13  Jun 08
Mac    CCS  3f086a62   [=======]
Mac    CCS  add0d69f          [===========]
Mac    CCS  ca6a2dde                  [========]
Mac    rgpu dd63080c                  [========]
Linux  CCS  d1bcaa09                     [=====]
Mac    rgpu e493c2fd                          [============]
Linux  rgpu daf9f59c                            |
Linux  rgpu b74f64ed                              [===]
Linux  rgpu e80de831                                 [===]
Linux  rgpu da4896fe                                    [====]
Linux  rgpu a9eacbbc                                       [=]
Linux  rgpu fd509bb7                                        [========]
Mac    rgpu a910ecaa  *                                       [======]
Linux  rgpu 4862cd6d  *                                          |
Mac    rgpu 11e6e374  *                                             [========================]
Linux  rgpu 61dc0523  *                                              |
Linux  rgpu afa1ce4c  *                                              [=======================]

  [ ]/| = active session     * = full JSONL on disk (others: history.jsonl prompts survive 30-day eviction)
  Max overlap: 26d on rust-gpu (afa1ce4c Linux  /  11e6e374 Mac, May 14 → Jun 8)
""", size=7)

# --- 7. Section break: CCS ------------------------------------------------
s = add(2)
set_text(get_ph(s, 0), "Case study 1\nCCS — directing the agent on familiar code")

# --- 8. CCS in 60 seconds -------------------------------------------------
s = add(4)
set_text(get_ph(s, 0), "CCS in sixty seconds")
set_text(get_ph(s, 13),
    "C99 library exposing autotuning configuration spaces, objective spaces, and tuners with an "
    "ask/tell pattern. Parameters can be numerical, categorical, ordinal, or string; spaces support "
    "conditions, forbidden clauses, expression trees, and feature contexts. Reference-counted objects, "
    "JSON and binary serialization, Python and Ruby bindings, Kokkos profiling connector. "
    "Used by autotuning workflows where multiple frameworks need to interoperate.")

# --- 9. CCS in code (showcase) --------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "CCS — what user code looks like")
set_text(get_ph(s, 13), "Python binding; equivalent C, Ruby APIs exist")
set_text(get_ph(s, 16), "Code")
set_code(get_ph(s, 14),
"""import cconfigspace as ccs

# 1. Define the search space
h1 = ccs.NumericalParameter.Float(lower=-5, upper=5)
h2 = ccs.NumericalParameter.Float(lower=-5, upper=5)
cs = ccs.ConfigurationSpace(parameters=[h1, h2])

# 2. Define the objective
v  = ccs.NumericalParameter.Float(...)
os = ccs.ObjectiveSpace(parameters=[v],
                        objectives={v: ccs.MINIMIZE})

# 3. Ask/tell loop
tuner = ccs.RandomTuner(objective_space=os)
for _ in range(100):
    configs, _ = tuner.ask(1)
    evals = [my_eval(c) for c in configs]
    tuner.tell(evals)

best = tuner.optima
""")
set_text(get_ph(s, 17), "What the agent worked on")
set_bullets(get_ph(s, 15), [
    "JSON serialization for every object class above",
    "Typed deserialize (no serialize-in-deserialize)",
    "Ruby + Python parity with the C API",
    "gcov-driven bug hunt across the C core",
])

# --- 10. CCS workflow -----------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "CCS workflow with the agent")
set_text(get_ph(s, 13), "Branch-per-task; PR-per-task; aggressive rebasing")
set_bullets(get_ph(s, 14), [
    "Every change opens a branch off devel and a single-topic PR",
    "Agent commits as bricevideau-ai, co-authored by Claude",
    ("I am the reviewer — every PR is read and merged manually", 1),
    "Rebase devel on upstream after merge — drilled in as a memory rule",
    "CI: clang-format-17, sanitizers, valgrind, multi-binding tests",
    ("Agent regressed to plain clang-format once — corrected, then memorized", 1),
    "make check / make check-valgrind is the contract for green",
])

# --- 11. CCS what landed --------------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "CCS — what landed in six weeks")
set_text(get_ph(s, 13), "96 PRs merged on argonne-lcf/CCS by the AI persona")
set_text(get_ph(s, 16), "Numbers")
set_bullets(get_ph(s, 14), [
    "99 PRs opened (#24 → #122)",
    "96 merged, 3 closed (~97% merge rate)",
    "~16 PRs per week sustained",
    "Each PR reviewable in <10 minutes",
])
set_text(get_ph(s, 17), "Topics by volume")
set_bullets(get_ph(s, 15), [
    "JSON serialization stack — extract/embed/array/scalar helpers",
    "Deserializer redesign — typed, no serialize-in-deserialize",
    "Bug fixes surfaced by gcov coverage analysis",
    "Binding parity — Ruby + Python exposure",
    "Doxygen + Markdown documentation rewrites",
])

# --- 12. CCS how I interacted (4-block) -----------------------------------
s = add(9)
set_text(get_ph(s, 0), "CCS — how I interacted")
set_text(get_ph(s, 13), "Short, directive prompts; frequent corrections")
set_block(get_ph(s, 16), "Direct correction",
    "“No, check again.” — said often, without preamble. The agent treated as a junior dev.")
set_block(get_ph(s, 17), "Catching toolchain drift",
    "“I noticed you started using clang-format instead of clang-format-17.”")
set_block(get_ph(s, 18), "Holding a quality line",
    "“I don’t think computer scientists should be lazy, it usually ends up costing us more in the long run.”")
set_block(get_ph(s, 19), "Project idioms",
    "“Could you please use CCS_REFUTE_ERR_GOTO instead?” Macros must be taught — then re-taught after each compaction.")

# --- 13. CCS wins ---------------------------------------------------------
s = add(6)
set_text(get_ph(s, 0), "CCS — wins")
set_text(get_ph(s, 13), "Where the agent unambiguously moved the project forward")
set_text(get_ph(s, 16), "Coverage push")
set_bullets(get_ph(s, 14), [
    "gcov tooling added as CI step",
    "Coverage gaps → test PRs",
    "Surfaced uninit-after-goto, realloc cleanup, deserialize bounds",
])
set_text(get_ph(s, 17), "Format and docs")
set_bullets(get_ph(s, 15), [
    "Full JSON serialization across every CCS object class",
    "Versioned format header, round-trip tests",
    "Deserialize made typed — code-smell eliminated",
])
set_text(get_ph(s, 19), "Binding parity")
set_bullets(get_ph(s, 18), [
    "Ruby + Python: missing C functions exposed",
    "Tests written for previously uncovered paths",
    "Critical binding bugs fixed",
])

# --- 14. CCS pitfalls -----------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "CCS — pitfalls observed")
set_text(get_ph(s, 13), "Even on owned code, the agent leaks bad habits unless gated")
set_bullets(get_ph(s, 14), [
    "Silent toolchain swaps (clang-format vs -17): caught by visual diff of CI output",
    "Default-to-laziness on cleanup paths and overflow checks",
    "Out-of-scope edits crept in: “you accidentally modified lines out of scope”",
    "Repeated reminders needed for macro/idiom families (CCS_REFUTE_ERR_GOTO)",
    "Doc format tensions: “the html is fixed but the markdown is now broken”",
    ("Mitigation: short PRs + active review + memory + CLAUDE.md", 1),
])

# --- 15. Section break: rust-gpu/claspr -----------------------------------
s = add(2)
set_text(get_ph(s, 0), "Case study 2\nrust-gpu + claspr — across a knowledge gap")

# --- 16. rust-gpu/claspr 60 sec -------------------------------------------
s = add(4)
set_text(get_ph(s, 0), "rust-gpu + claspr in sixty seconds")
set_text(get_ph(s, 13),
    "rust-gpu compiles Rust to SPIR-V; it targeted Vulkan only. We added the OpenCL Kernel execution model: "
    "Physical64 addressing, mandatory OpenCL capabilities, slice decomposition into (ptr, len) pairs, "
    "OpenCL.std math intrinsics, native CL vector types (Float4, Int8, Long16, swizzles), DebugPrintf, "
    "subgroup and work-group collectives, storage images and samplers, atomic flags, f64 throughout. "
    "claspr is the host-side layer on top: a #[claspr::device] proc-macro extracts a kernel sub-crate at "
    "build time, generates a typed launch wrapper, and gives users single-source OpenCL in Rust.")

# --- 17. claspr in code (showcase) ----------------------------------------
s = add(5)
set_text(get_ph(s, 0), "claspr — what user code looks like")
set_text(get_ph(s, 13), "One Rust file: device kernel + host driver, side by side")
set_text(get_ph(s, 16), "Code")
set_code(get_ph(s, 14),
"""use claspr::{Context, DeviceSlice};

#[claspr::device]
mod gpu {
    #[claspr::kernel]
    pub fn collatz_kernel(
        #[spirv(global_invocation_id)] id: glam::USizeVec3,
        #[spirv(cross_workgroup)] data: &mut [u32],
    ) {
        data[id.x] = collatz(data[id.x]);
    }
}

fn main() -> claspr::Result<()> {
    let ctx = Context::any()?;
    let kernels = gpu::kernels(&ctx)?;
    let mut h: Vec<u32> = (1..=1024).collect();
    let mut d = DeviceSlice::<u32>::alloc_zero(&ctx, h.len())?;
    d.write(&h).wait(&ctx)?;
    let d = kernels.collatz_kernel([1024], d).wait(&ctx)?;
    d.read(&mut h).wait(&ctx)?;
    Ok(())
}
""")
set_text(get_ph(s, 17), "What the macros + chain do")
set_bullets(get_ph(s, 15), [
    "#[claspr::device] mod gpu: lifted into its own SPIR-V crate at build time",
    "#[claspr::kernel]: marks entry point; emits typed launch wrapper",
    ".write(&h).wait(&ctx)?: Tier 1 op + terminal wait; Tier 2 chains use .and_then(...)",
    "Same file built twice: host (cargo) + SPIR-V (rust-gpu via build.rs)",
])

# --- 18. Knowledge-gap shift (4-block) ------------------------------------
s = add(9)
set_text(get_ph(s, 0), "The knowledge-gap shift")
set_text(get_ph(s, 13), "Same engineer, same tool — completely different mode")
set_block(get_ph(s, 16), "What I do know",
    "OpenCL spec, SPIR-V semantics, HPC kernel perf-eng, library and runtime design, implementation quirks.")
set_block(get_ph(s, 17), "What I do not know",
    "Rust idioms, rust-gpu codegen internals, proc-macro authoring, async Rust, Rust GPU conventions.")
set_block(get_ph(s, 18), "What this changes",
    "I cannot pre-write the answer in my head. I have to let the agent enumerate options and reason out loud.")
set_block(get_ph(s, 19), "New rhythm",
    "Explore → enumerate → discuss → decide → implement → review. Agent enumerates; engineer judges.")

# --- 19. rust-gpu what landed ---------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "rust-gpu + claspr — what landed")
set_text(get_ph(s, 13), "One large fork PR + one new repository + a full test infrastructure")
set_text(get_ph(s, 16), "rust-gpu PR #3 (OpenCL Kernel)")
set_bullets(get_ph(s, 14), [
    "+24,540 / -595 lines across 761 files",
    "Eight spirv-unknown-opencl* targets",
    "OpenCL Kernel execution model end-to-end",
    "spirv-std: cl::*, opencl_std, printf!, atomics, groups",
    "Difftests on rusticl/llvmpipe in CI",
])
set_text(get_ph(s, 17), "claspr (greenfield)")
set_bullets(get_ph(s, 15), [
    "182 commits, ~81K LOC of Rust",
    "Proc-macro single-source: #[claspr::device], #[claspr::kernel]",
    "Two-tier op model: Tier 1 ops + Tier 2 async chains",
    "Access markers (ReadWrite/ReadOnly/Frozen/…)",
    "Examples: raymarch, mandelbrot, sobel, image-pipeline",
])

# --- 20. rust-gpu workflow ------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "rust-gpu / claspr workflow")
set_text(get_ph(s, 13), "Branch discipline, cross-laptop handoff, upstream-aware")
set_bullets(get_ph(s, 14), [
    "Started from upstream issue #74 — reproduce, then design",
    "Two long-lived branches: opencl-kernel-support (stable) and -v2 (dev)",
    ("Samples repo tracks stable; promotion is explicit, not automatic", 1),
    "Cross-laptop hand-offs between Intel-Ubuntu and Arm-Mac sessions",
    "Lint-before-push enforced after one CI failure (memorized)",
    "Rebase after merge enforced after the agent forgot once (memorized)",
    "claspr CI runs on rusticl/llvmpipe so reviewers can rerun",
])

# --- 21. rust-gpu how I interacted (4-block) ------------------------------
s = add(9)
set_text(get_ph(s, 0), "rust-gpu — interaction in knowledge-gap mode")
set_text(get_ph(s, 13), "Directive shrinks, dialogue grows")
set_block(get_ph(s, 16), "Scope-cutting",
    "“Alright, let me help refocus. I would be fine with a test that leverages the spirv(kernel) model, not a crate.”")
set_block(get_ph(s, 17), "Memory-anchored",
    "“I remember you telling me that the runner put the arguments in different order than the sources.”")
set_block(get_ph(s, 18), "Collaborative design",
    "“I want to have a conversation around Async/Sync and InOrder/OutOfOrder. But first — can Rust polymorphism depend on the return value?”")
set_block(get_ph(s, 19), "Pre-emptive context discipline",
    "“You can create a WIP document describing the design as we progress — never know when compaction will bite us.”")

# --- 22. Design dialogues -------------------------------------------------
s = add(6)
set_text(get_ph(s, 0), "Design dialogues that produced real decisions")
set_text(get_ph(s, 13), "Cases where the agent enumerated trade-offs and I picked")
set_text(get_ph(s, 16), "Execution model")
set_bullets(get_ph(s, 14), [
    "Async vs sync, in-order vs out-of-order",
    "and_then chain returning values",
    "with_context as access-only, not sync escape",
    "Default in-order per device; OOO opt-in",
])
set_text(get_ph(s, 17), "Memory and safety")
set_bullets(get_ph(s, 15), [
    "alloc_uninit marked unsafe (no MaybeUninit in rust-gpu)",
    "Access markers as type-state",
    "HostBuffer persistent-map flagged UB per spec",
    "Arc<DeviceSlice> read-only, write-then-share",
])
set_text(get_ph(s, 19), "Abstractions")
set_bullets(get_ph(s, 18), [
    "Tier 1 (ops) decoupled from Tier 2 (async)",
    "KernelOp trait so proc-macro stays light",
    "Late-bind launcher: buf.write(&data).wait(&ctx)?",
    "Spike scenarios kept as design docs",
])

# --- 23. rust-gpu pitfalls ------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "rust-gpu / claspr — pitfalls observed")
set_text(get_ph(s, 13), "Different failure modes when expertise is asymmetric")
set_bullets(get_ph(s, 14), [
    "Context exhaustion: 8 auto-compactions in one rolling session",
    "Drift after compaction: required explicit difftest run-command memory",
    "Cleaning up legitimate spike code as if it were dead code",
    "Tool/version regressions: pocl rebuilds + LLVM drift → silent kernel hangs",
    "Glam transitive-dep mismatch: “I don’t understand how this happened”",
    ("Mitigation: end sessions; preempt compaction; pin versions; memory entries", 1),
])

# --- 24. Section break: cross-cutting -------------------------------------
s = add(2)
set_text(get_ph(s, 0), "Cross-cutting analysis\nWhat the data shows across both projects")

# --- 25. Tool-call distribution -------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Tool-call distribution")
set_text(get_ph(s, 13), "From the 5 locally-available rust-gpu sessions (~12,000 calls)")
set_text(get_ph(s, 16), "Top tools by call count")
set_bullets(get_ph(s, 14), [
    "Bash         6,357   ~53%",
    "Edit         2,291   ~19%",
    "Read         1,342   ~11%",
    "TaskUpdate     636   (todo tracking)",
    "Write          548   (new files)",
    "TaskCreate, Monitor, Agent, WebFetch: long tail",
])
set_text(get_ph(s, 17), "What this means")
set_bullets(get_ph(s, 15), [
    "Bash dominates — agent lives in the shell, not just the editor",
    "~1.7 Edits per Read — agent reads a chunk, then makes several small targeted changes within it",
    "TaskUpdate + TaskCreate ~8% — todo lists are load-bearing UX",
    "Write small — most output goes into existing files",
])

# --- 26. Sub-agent adoption curve -----------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Sub-agent adoption curve")
set_text(get_ph(s, 13), "Adoption is workflow-driven, not just model-driven")
set_text(get_ph(s, 16), "By surviving session (rust-gpu)")
set_bullets(get_ph(s, 14), [
    "Mac    a910ecaa (May 4–13):   0",
    "Linux  4862cd6d (May 11):      0",
    "Linux  61dc0523 (May 14):      0",
    "Mac    11e6e374 (May 13–Jun 8):  24",
    "Linux  afa1ce4c (May 14–Jun 8):  1",
    "Same engineer, same model, 26 parallel days: 24×",
])
set_text(get_ph(s, 17), "What this means")
set_bullets(get_ph(s, 15), [
    "Model era enables sub-agents; workflow choice gates use",
    "Mac: design-dialogue on claspr — Explore fan-out paid off",
    "Linux: capability-audit + spec deep-dives — fewer branches to map",
    "Caveat: 30-day eviction hides Mar–Apr sessions",
])

# --- 27. Model progression ------------------------------------------------
# The model itself was a moving target across the work — likely explains
# some of the tooling-adoption drift on the prior slide.
s = add(5)
set_text(get_ph(s, 0), "Model progression during the work")
set_text(get_ph(s, 13), "The model itself was a moving target; usage style moved with it")
set_text(get_ph(s, 16), "Timeline (from commit metadata)")
set_bullets(get_ph(s, 14), [
    "Feb 25–27: Sonnet 4.6 — first days on CCS (7 commits)",
    "Feb 27 → Apr 7: Opus 4.6 carries most of CCS (232 commits)",
    "Mar 27: Opus 4.6 1M-context arrives mid-CCS",
    "Apr 24: Opus 4.7 arrives mid-rust-gpu",
    "claspr (May–Jun): entirely Opus 4.7, often 1M",
])
set_text(get_ph(s, 17), "What tracked the upgrades")
set_bullets(get_ph(s, 15), [
    "Observed mid-stream: “your model was updated since we started on this project” (Apr 22)",
    "1M context enabled longer design dialogues — coincides with shift to dialogue mode",
    "Sub-agent adoption emerges only in the Opus 4.7 era",
    "More willing to enumerate trade-offs vs pick a default",
])

# --- 28. Token economics (4-block) ----------------------------------------
s = add(9)
set_text(get_ph(s, 0), "Token economics")
set_text(get_ph(s, 13), "ccusage covers 36% of prompts directly; the rest is projection")
set_block(get_ph(s, 16), "Measured (May–Jun, both laptops)",
    "$3,474 across 5 surviving rust-gpu sessions / 963 prompts. ≈ $3.61/prompt at long-dialogue rates.")
set_block(get_ph(s, 17), "Projected total",
    "Scaling to all 2,659 prompts: ≈ $8.6–9.6 K for the full 15 weeks. CCS prompts were terser, likely cheaper per turn.")
set_block(get_ph(s, 18), "April cache bug",
    "An April Claude Code caching bug inflated the early-rust-gpu sessions; magnitude lost with evicted JSONLs. Treat projection as a floor.")
set_block(get_ph(s, 19), "Volume (visible portion)",
    "9.5 B cache reads vs 334 M cache writes vs 431 K raw input vs 18.7 M output. Cache is the workhorse.")

# --- 29. Memory system ----------------------------------------------------
s = add(6)
set_text(get_ph(s, 0), "Memory system in practice")
set_text(get_ph(s, 13), "30+ entries built up; indexed by MEMORY.md")
set_text(get_ph(s, 16), "feedback (rules)")
set_bullets(get_ph(s, 14), [
    "clang-format-17, not clang-format",
    "Rebase devel after merge, unprompted",
    "cargo fmt + clippy before push",
    "Respect spike intent",
])
set_text(get_ph(s, 17), "project (decisions)")
set_bullets(get_ph(s, 15), [
    "OpenCL kernel stable promoted from -v2",
    "Access markers landed (5 markers + tests)",
    "alloc_uninit marked unsafe",
    "with_context removed",
])
set_text(get_ph(s, 19), "reference (where to look)")
set_bullets(get_ph(s, 18), [
    "pocl 7.2-pre at ~/.local (patched)",
    "Intel OpenCL Intercept Layer",
    "macOS opencl3 picks Apple framework",
    "Rusticl queue unusable after termination",
])

# --- 30. Context-window management ----------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Context-window management is part of the workflow")
set_text(get_ph(s, 13), "200K is large until it isn’t")
set_text(get_ph(s, 16), "Symptoms")
set_bullets(get_ph(s, 14), [
    "Auto-compaction: 8 distinct days in one session",
    "Each compaction summarizes prior turns",
    "Quality dips: forgotten conventions",
    "Memory entries are the only ground truth across compactions",
])
set_text(get_ph(s, 17), "Practices that worked")
set_bullets(get_ph(s, 15), [
    "Preempt: “be careful context is almost full, please anticipate”",
    "Manual /compact and /clear",
    "WIP design docs written to disk",
    "End-of-session: update CLAUDE.md and memory",
    "Sub-agents absorb cost instead of inflating context",
])

# --- 31. Evolution CCS → claspr -------------------------------------------
s = add(5)
set_text(get_ph(s, 0), "Evolution of interaction style")
set_text(get_ph(s, 13), "From short directives to long design dialogues")
set_text(get_ph(s, 16), "Early (CCS era, Feb–Apr)")
set_bullets(get_ph(s, 14), [
    "Short prompts, often <30 words",
    "PR-per-task, ship-fast cadence",
    "Zero sub-agents, zero skills, ~no slash cmds",
    "Memory hand-rolled via README + CLAUDE.md",
    "Errors caught in human review, not agent self-review",
])
set_text(get_ph(s, 17), "Late (claspr era, May–Jun)")
set_bullets(get_ph(s, 15), [
    "Longer prompts with design questions",
    "Explore sub-agents fan out research in parallel",
    "Plan sub-agent before non-trivial implementation",
    "30+ memory entries; CLAUDE.md as orientation contract",
    "Pre-emptive /compact; agent enumerates, engineer picks",
])

# --- 32. What worked best (4-block) ---------------------------------------
s = add(9)
set_text(get_ph(s, 0), "What worked best")
set_text(get_ph(s, 13), "Patterns the data and the experience agree on")
set_block(get_ph(s, 16), "Expert + agent on familiar code",
    "CCS: 16 PRs/week sustained, 97% merge rate. Expert pays review tax; agent pays typing tax.")
set_block(get_ph(s, 17), "Explore-then-decide on new code",
    "Parallel Explore sub-agents survey the design space; engineer makes the call.")
set_block(get_ph(s, 18), "Short PRs + active review",
    "Single-topic, mergeable in ~10 min review. Out-of-scope creep caught immediately.")
set_block(get_ph(s, 19), "Memory + CLAUDE.md as ground truth",
    "Survives compaction, sessions, laptops. Conventions encoded once, not re-litigated.")

# --- 33. Anti-patterns (4-block) ------------------------------------------
s = add(9)
set_text(get_ph(s, 0), "Anti-patterns observed")
set_text(get_ph(s, 13), "Failure modes seen at least once each")
set_block(get_ph(s, 16), "One rolling forever-session",
    "Hot state is tempting. But every compaction is a quality hit. End sessions explicitly.")
set_block(get_ph(s, 17), "Trusting CI-green without diff",
    "Toolchain swaps, dep drift, out-of-scope edits all pass CI. Diff review is not optional.")
set_block(get_ph(s, 18), "Letting the agent pick the toolchain",
    "clang-format vs -17, glam version, pocl local vs system. Pin versions in CLAUDE.md.")
set_block(get_ph(s, 19), "Cleaning up unfamiliar patterns",
    "Agent removes spike code or workarounds it doesn’t recognize. Document intent inline first.")

# --- 34. Implications -----------------------------------------------------
s = add(6)
set_text(get_ph(s, 0), "Implications for performance engineering")
set_text(get_ph(s, 13), "Where AI agents fit and where they don’t")
set_text(get_ph(s, 16), "Net-positive cases")
set_bullets(get_ph(s, 14), [
    "Mature codebases the engineer owns",
    "Coverage push, doc rewrite, binding parity",
    "Repetitive fix campaigns from CI signals",
    "Multi-language interop work",
])
set_text(get_ph(s, 17), "Conditional cases")
set_bullets(get_ph(s, 15), [
    "Unfamiliar code with clear spec",
    "Greenfield with strong design dialogue",
    "Cross-language interop, with careful review",
])
set_text(get_ph(s, 19), "Avoid for now")
set_bullets(get_ph(s, 18), [
    "Microbenchmark optimization without validation",
    "Codebases where the spec is the code itself",
    "Time-pressured work with no review headroom",
    "Single rolling session beyond a few days",
])

# --- 35. Recommendations --------------------------------------------------
s = add(3)
set_text(get_ph(s, 0), "Recommendations — starter checklist")
set_text(get_ph(s, 13), "If you start an agent-assisted project tomorrow")
set_bullets(get_ph(s, 14), [
    "Give the agent its own GitHub identity",
    "Write a CLAUDE.md from day 1 (conventions, build, test, taboos)",
    "Maintain memory: feedback rules, project decisions, references",
    "Short PRs, one topic, mergeable in <10 min review",
    "Pin toolchain versions explicitly",
    "End sessions before they end you; preempt /compact",
    "Use Explore sub-agents whenever you don’t know the answer",
    "Trust diff inspection, not CI-green, as the merge gate",
])

# --- 36. Closing ----------------------------------------------------------
add(13)

pres.save(OUT)
print(f"Wrote {OUT}")
print(f"Slide count: {len(pres.slides)}")
